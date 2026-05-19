import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Emu
from collections import Counter

prs = Presentation('C:/GitHub/Kyn-skyn/output/Kyn_skyn_updated.pptx')
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

print(f'Slide dimensions: {SLIDE_W/914400:.2f}" x {SLIDE_H/914400:.2f}"')
print(f'Total slides: {len(prs.slides)}')
print('=' * 70)

total_issues = 0

for i, slide in enumerate(prs.slides):
    issues = []
    shape_details = []

    for shape in slide.shapes:
        left, top = shape.left, shape.top
        w, h = shape.width, shape.height
        txt = ''
        if hasattr(shape, 'text'):
            txt = shape.text.strip()

        shape_details.append({
            'name': shape.name,
            'left': left, 'top': top, 'width': w, 'height': h,
            'text': txt,
            'shape_type': str(shape.shape_type),
        })

        # Off-slide elements (with tolerance)
        if left + w > SLIDE_W + 50000:
            issues.append(f'OFF-SLIDE RIGHT: "{shape.name}" extends to {(left+w)/914400:.2f}" (slide width {SLIDE_W/914400:.2f}")')
        if top + h > SLIDE_H + 50000:
            issues.append(f'OFF-SLIDE BOTTOM: "{shape.name}" extends to {(top+h)/914400:.2f}" (slide height {SLIDE_H/914400:.2f}")')
        if left < -10000:
            issues.append(f'OFF-SLIDE LEFT: "{shape.name}" at x={left/914400:.2f}"')
        if top < -10000:
            issues.append(f'OFF-SLIDE TOP: "{shape.name}" at y={top/914400:.2f}"')

        # Zero-size elements with text
        if (w == 0 or h == 0) and txt:
            issues.append(f"ZERO-SIZE WITH TEXT: \"{shape.name}\" ({w}x{h}) text='{txt[:40]}'")

        # Very small elements
        if w < 5000 and h < 5000 and w > 0 and h > 0:
            issues.append(f"TINY ELEMENT: \"{shape.name}\" ({w/914400:.4f}\" x {h/914400:.4f}\") text='{txt[:20]}'")

        # Elements at exactly (0,0)
        if left == 0 and top == 0:
            issues.append(f'AT ORIGIN (0,0): "{shape.name}" size=({w/914400:.2f}" x {h/914400:.2f}") text="{txt[:40]}"')

        # Empty text boxes (shapes with text frame but no content)
        if hasattr(shape, 'text_frame') and not txt and shape.has_text_frame:
            # Only flag if it looks like it should have text (not decorative shapes)
            if 'text' in shape.name.lower() or 'title' in shape.name.lower() or 'subtitle' in shape.name.lower():
                issues.append(f'EMPTY TEXT ELEMENT: "{shape.name}" at ({left/914400:.2f}", {top/914400:.2f}")')

        # Very large shapes that might be covering content
        if w > SLIDE_W * 1.5 or h > SLIDE_H * 1.5:
            issues.append(f'OVERSIZED: "{shape.name}" is {w/914400:.2f}" x {h/914400:.2f}" (larger than 1.5x slide)')

    # Check for duplicate shape names with detail
    names = [s.name for s in slide.shapes]
    name_counts = Counter(names)
    for name, count in name_counts.items():
        if count > 1:
            dupes = [sd for sd in shape_details if sd['name'] == name]
            issues.append(f'DUPLICATE NAME (x{count}): "{name}"')
            for d in dupes:
                pos_info = f'  -> pos=({d["left"]/914400:.2f}", {d["top"]/914400:.2f}") size=({d["width"]/914400:.2f}" x {d["height"]/914400:.2f}")'
                if d['text']:
                    pos_info += f' text="{d["text"][:50]}"'
                issues.append(pos_info)

    # Check for shapes at identical positions (potential overlaps/ghosts)
    positions = {}
    for sd in shape_details:
        key = (sd['left'], sd['top'], sd['width'], sd['height'])
        if key in positions:
            issues.append(f'EXACT OVERLAP: "{sd["name"]}" and "{positions[key]}" at ({sd["left"]/914400:.2f}", {sd["top"]/914400:.2f}") size ({sd["width"]/914400:.2f}" x {sd["height"]/914400:.2f}")')
        positions[key] = sd['name']

    if issues:
        print(f'\nSLIDE {i+1} ({len(slide.shapes)} shapes):')
        for issue in issues:
            print(f'  {issue}')
        total_issues += len(issues)

print(f'\n{"=" * 70}')
print(f'TOTAL ISSUES FOUND: {total_issues}')
print(f'\n=== SUMMARY BY SLIDE ===')
for i, slide in enumerate(prs.slides):
    print(f'  Slide {i+1:2d}: {len(slide.shapes):3d} shapes')
print(f'\n=== CHECK COMPLETE ===')
