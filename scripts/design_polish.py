"""
Add design polish to slides 3, 4, 6, 7:
- Outer shadows on all containers
- Rounded corners
- Left-accent borders
- Callout container styling
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

prs = Presentation('C:/GitHub/Kyn-skyn/output/Kyn_skyn_final.pptx')

def add_shadow(shape, blur=50800, dist=38100, direction=2700000, alpha_pct=15):
    """Add outer shadow matching original deck cards."""
    spPr = shape._element.spPr
    for existing in spPr.findall(qn('a:effectLst')):
        spPr.remove(existing)
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    shdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    shdw.set('blurRad', str(blur))
    shdw.set('dist', str(dist))
    shdw.set('dir', str(direction))
    shdw.set('algn', 'bl')
    shdw.set('rotWithShape', '0')
    clr = etree.SubElement(shdw, qn('a:srgbClr'))
    clr.set('val', '000000')
    a = etree.SubElement(clr, qn('a:alpha'))
    a.set('val', str(alpha_pct * 1000))

def set_rounded(shape, adj=5000):
    """Set rounded corners on a shape."""
    spPr = shape._element.spPr
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')
        for avLst in prstGeom.findall(qn('a:avLst')):
            prstGeom.remove(avLst)
        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {adj}')

def add_border(shape, color_hex, width=12700):
    """Add/replace border on shape."""
    spPr = shape._element.spPr
    for ln in spPr.findall(qn('a:ln')):
        spPr.remove(ln)
    ln = etree.SubElement(spPr, qn('a:ln'))
    ln.set('w', str(width))
    sf = etree.SubElement(ln, qn('a:solidFill'))
    c = etree.SubElement(sf, qn('a:srgbClr'))
    c.set('val', color_hex)

# ================================================================
# SLIDE 3 (Why Now) - Add shadows, rounded corners, borders to cards
# ================================================================
print("Polishing Slide 3 (Why Now)...")
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if shape.shape_type is None:
        continue
    try:
        stype = str(shape.shape_type)
    except:
        continue

    # Card backgrounds (rectangles with colored fills)
    if 'RECTANGLE' in stype and not 'ROUNDED' in stype:
        # Check if it's a card (not the accent bar, not header)
        w = shape.width
        h = shape.height
        if w > Emu(3000000) and h > Emu(800000) and h < Emu(2000000):
            # This is a content card
            set_rounded(shape, adj=5000)
            add_shadow(shape, blur=50800, dist=38100, alpha_pct=15)
            add_border(shape, 'D0D0D0', 12700)
            print(f"  Card: {shape.name} - shadow + rounded + border")

        # Bottom callout bar
        if w > Emu(8000000) and h > Emu(300000) and h < Emu(600000) and shape.top > Emu(5000000):
            set_rounded(shape, adj=8000)
            add_shadow(shape, blur=38100, dist=25400, alpha_pct=10)
            print(f"  Callout bar: {shape.name} - shadow + rounded")

    # Number badges - add shadow
    if 'OVAL' in stype:
        if shape.width < Emu(500000):
            add_shadow(shape, blur=25400, dist=12700, alpha_pct=20)
            print(f"  Badge: {shape.name} - shadow")

print("  Done")

# ================================================================
# SLIDE 4 (Investable) - Add shadows, left-accent borders, quote styling
# ================================================================
print("Polishing Slide 4 (Investable)...")
slide4 = prs.slides[3]

for shape in slide4.shapes:
    try:
        stype = str(shape.shape_type)
    except:
        continue

    if 'RECTANGLE' not in stype:
        continue

    w = shape.width
    h = shape.height

    # Header bars (dark colored, full width, ~400000 tall)
    if w > Emu(8000000) and h > Emu(300000) and h < Emu(500000):
        set_rounded(shape, adj=3000)
        add_shadow(shape, blur=25400, dist=12700, alpha_pct=12)
        print(f"  Header bar: {shape.name} ({shape.top}) - rounded + shadow")

    # Pre-pitch validation box (large amber area at bottom)
    if w > Emu(8000000) and h > Emu(1500000):
        set_rounded(shape, adj=4000)
        add_shadow(shape, blur=50800, dist=38100, alpha_pct=15)
        add_border(shape, 'E8C547', 19050)  # 1.5pt gold border
        print(f"  Validation box: {shape.name} - rounded + shadow + gold border")

    # Right-edge accent bar (tall, narrow, red)
    if w < Emu(600000) and h > Emu(6000000):
        # This is the accent bar - no changes needed
        pass

# Add left-accent bars to the thesis and investor sections
# Find the header bar positions to know where to place accents
for shape in slide4.shapes:
    if hasattr(shape, 'text') and shape.text.strip() == 'THE THESIS':
        # Add a left accent bar for the thesis section
        thesis_top = shape.top
    if hasattr(shape, 'text') and shape.text.strip() == 'INVESTOR FIT':
        investor_top = shape.top

# Thesis left accent - navy
bar1 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Emu(650000), Emu(1200000), Emu(50000), Emu(1500000))
bar1.fill.solid()
bar1.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x61)
bar1.line.fill.background()
print("  Added navy left-accent for THESIS")

# Investor left accent - teal
bar2 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Emu(650000), Emu(3000000), Emu(50000), Emu(1500000))
bar2.fill.solid()
bar2.fill.fore_color.rgb = RGBColor(0x0D, 0x73, 0x77)
bar2.line.fill.background()
print("  Added teal left-accent for INVESTOR FIT")

print("  Done")

# ================================================================
# SLIDE 6 (Flywheel) - Shadows, rounded corners, borders on nodes
# ================================================================
print("Polishing Slide 6 (Flywheel)...")
slide6 = prs.slides[5]

for shape in slide6.shapes:
    try:
        stype = str(shape.shape_type)
    except:
        continue

    if 'ROUNDED' in stype or 'RECTANGLE' in stype:
        w = shape.width
        h = shape.height

        # Flywheel node boxes (2400000 x 1200000)
        if w > Emu(2000000) and w < Emu(3000000) and h > Emu(1000000) and h < Emu(1500000):
            set_rounded(shape, adj=8000)
            add_shadow(shape, blur=63500, dist=38100, alpha_pct=18)
            print(f"  Node: {shape.name} - shadow + rounded")

    # Center oval - add shadow
    if 'OVAL' in stype:
        if shape.width > Emu(800000):
            add_shadow(shape, blur=50800, dist=25400, alpha_pct=20)
            print(f"  Center oval: {shape.name} - shadow")

print("  Done")

# ================================================================
# SLIDE 7 (Curology) - Shadow on KYN SKYN box, rounded corners, callout
# ================================================================
print("Polishing Slide 7 (Curology)...")
slide7 = prs.slides[6]

for shape in slide7.shapes:
    try:
        stype = str(shape.shape_type)
    except:
        continue

    if 'ROUNDED' not in stype and 'RECTANGLE' not in stype:
        continue

    w = shape.width
    h = shape.height

    # Comparison boxes (large, tall)
    if w > Emu(4000000) and h > Emu(2500000):
        set_rounded(shape, adj=5000)

        # Check if this is the Kyn Skyn box (right side, x > 5000000)
        if shape.left > Emu(5000000):
            # KYN SKYN box - stronger shadow for visual hierarchy
            add_shadow(shape, blur=63500, dist=38100, alpha_pct=20)
            print(f"  KYN SKYN box: {shape.name} - strong shadow + rounded")

            # Add gold left-accent bar
            accent = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                shape.left, shape.top, Emu(57150), shape.height)
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor(0xC5, 0x44, 0x27)  # Brand red
            accent.line.fill.background()
            print("  Added red left-accent to KYN SKYN box")
        else:
            # Curology box - subtle shadow
            add_shadow(shape, blur=38100, dist=25400, alpha_pct=10)
            print(f"  Curology box: {shape.name} - subtle shadow + rounded")

# ================================================================
# SAVE
# ================================================================
prs.save('C:/GitHub/Kyn-skyn/output/Kyn_skyn_final.pptx')
print("\nDesign polish complete.")
