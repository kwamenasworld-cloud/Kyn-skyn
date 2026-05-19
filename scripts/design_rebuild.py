"""
Rebuild slides 3, 4, 6, 7 using the original deck's design language.
Fix all footnotes to 8pt minimum.

Brand palette (from original slides):
- C54427: Brand red/rust (numbers, accents, callout bars)
- FEF2F2: Light blush (background panels)
- FFF4CE: Warm cream (product backgrounds)
- 0D7377: Dark teal (data bars)
- C45B28: Warm rust (secondary bars)
- 374151: Dark gray (body text)
- 6B7280: Medium gray (secondary text)
- 055528: Dark green (solution headers)
- F59E0B: Amber (warning/highlight)
- 1E2761: Deep navy (titles)
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

prs = Presentation('C:/GitHub/Kyn-skyn/output/Kyn_skyn_final.pptx')
ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

# Brand colors
RED = RGBColor(0xC5, 0x44, 0x27)
BLUSH = RGBColor(0xFE, 0xF2, 0xF2)
CREAM = RGBColor(0xFF, 0xF4, 0xCE)
TEAL = RGBColor(0x0D, 0x73, 0x77)
RUST = RGBColor(0xC4, 0x5B, 0x28)
BODY = RGBColor(0x37, 0x41, 0x51)
GRAY = RGBColor(0x6B, 0x72, 0x80)
DKGRAY = RGBColor(0x4B, 0x55, 0x63)
GREEN = RGBColor(0x05, 0x55, 0x28)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
NAVY = RGBColor(0x1E, 0x27, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def clear_slide(slide):
    """Remove all shapes from a slide except the required group properties."""
    spTree = slide.shapes._spTree
    to_remove = []
    for child in list(spTree):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'pic', 'grpSp', 'cxnSp', 'graphicFrame'):
            to_remove.append(child)
    for el in to_remove:
        spTree.remove(el)

def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, fill_rgb):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape

def add_tb(slide, left, top, width, height, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = wrap
    return tb

def set_run(run, text, size_pt, color, bold=False):
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold

def add_para(tf, text, size_pt, color, bold=False, space_before=None):
    p = tf.add_paragraph()
    if space_before:
        p.space_before = Pt(space_before)
    r = p.add_run()
    set_run(r, text, size_pt, color, bold)
    return p

# ================================================================
# REBUILD SLIDE 3 (Why Now) - using brand colors, sharp corners
# ================================================================
print("Rebuilding Slide 3 (Why Now)...")
slide3 = prs.slides[2]
clear_slide(slide3)
slide3.background.fill.solid()
slide3.background.fill.fore_color.rgb = WHITE

# Title using "Topic | Statement" format
tb = add_tb(slide3, Emu(850168), Emu(400000), Emu(10500000), Emu(550000))
p = tb.text_frame.paragraphs[0]
r1 = p.add_run(); set_run(r1, "Why now | ", 20, NAVY, True)
r2 = p.add_run(); set_run(r2, "Four forces converging to create this category", 20, BODY, True)

# Right-edge rust accent bar (matching original slides)
add_rect(slide3, Emu(11700000), Emu(0), Emu(492000), Emu(6858000), RED)

# Four numbered items using original deck's number-badge pattern
items = [
    ("Scalp care is surging", "Global scalp care market growing 8%+ CAGR. Mintel reports scalp health is the #1 emerging concern in skincare."),
    ("Social media is driving awareness", "TikTok #sebderm has 500M+ views. Reddit r/SebDerm has 80K+ members. Consumers are self-diagnosing but finding nothing purpose-built."),
    ("Telehealth is normalized", "Post-COVID adoption makes a $30 consult routine. Ola Digital Health offers turnkey setup for $2K. The infrastructure exists."),
    ("No incumbent owns this", "Curology serves acne. Dermazen sells products. Head & Shoulders treats symptoms. Nobody combines diagnosis + treatment for Malassezia."),
]

y_start = Emu(1200000)
row_h = Emu(1150000)
col_w = Emu(5100000)

for i, (title, body) in enumerate(items):
    row = i // 2
    col = i % 2
    x = Emu(700000) + col * (col_w + Emu(400000))
    y = y_start + row * (row_h + Emu(200000))

    # Light background strip
    bg_color = BLUSH if i % 2 == 0 else CREAM
    add_rect(slide3, x, y, col_w, row_h, bg_color)

    # Number badge (brand red circle)
    badge = add_circle(slide3, x + Emu(120000), y + Emu(100000), Emu(330000), RED)
    num_tb = add_tb(slide3, x + Emu(120000), y + Emu(100000), Emu(330000), Emu(330000))
    num_tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = num_tb.text_frame.paragraphs[0].add_run()
    set_run(r, str(i + 1), 14, WHITE, True)

    # Title text
    title_tb = add_tb(slide3, x + Emu(550000), y + Emu(100000), col_w - Emu(700000), Emu(330000))
    r = title_tb.text_frame.paragraphs[0].add_run()
    set_run(r, title, 16, BODY, True)

    # Body text
    body_tb = add_tb(slide3, x + Emu(120000), y + Emu(500000), col_w - Emu(240000), Emu(600000))
    r = body_tb.text_frame.paragraphs[0].add_run()
    set_run(r, body, 11, DKGRAY)

# Bottom callout using amber warning-bar style (from slide 5)
bar = add_rect(slide3, Emu(700000), Emu(5700000), Emu(10500000), Emu(450000), RGBColor(0xFE, 0xF3, 0xC7))
tb = add_tb(slide3, Emu(800000), Emu(5720000), Emu(10300000), Emu(400000))
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r1 = p.add_run(); set_run(r1, "The category is forming now. ", 14, AMBER, True)
r2 = p.add_run(); set_run(r2, "The question is who builds the platform first.", 14, BODY, False)


# ================================================================
# REBUILD SLIDE 4 (Why This Is Investable)
# ================================================================
print("Rebuilding Slide 4 (Investable)...")
slide4 = prs.slides[3]
clear_slide(slide4)
slide4.background.fill.solid()
slide4.background.fill.fore_color.rgb = WHITE

# Title
tb = add_tb(slide4, Emu(850168), Emu(400000), Emu(10500000), Emu(550000))
p = tb.text_frame.paragraphs[0]
r1 = p.add_run(); set_run(r1, "Why this is investable | ", 20, NAVY, True)
r2 = p.add_run(); set_run(r2, "Vertical integration with a clear validation path", 20, BODY, True)

# Right-edge rust accent bar
add_rect(slide4, Emu(11700000), Emu(0), Emu(492000), Emu(6858000), RED)

# Section 1: THE THESIS - navy header bar + content
add_rect(slide4, Emu(700000), Emu(1200000), Emu(10500000), Emu(400000), NAVY)
tb = add_tb(slide4, Emu(800000), Emu(1210000), Emu(10300000), Emu(380000))
r = tb.text_frame.paragraphs[0].add_run()
set_run(r, "THE THESIS", 16, WHITE, True)

thesis_points = [
    "Telehealth + product loop creates a data moat no competitor can replicate",
    "Clean unit economics: $14 contribution margin, 75%+ product margin",
    "$250K ask on $1.5M cap — 18 months to breakeven",
    "Not just skincare, not just telehealth — vertical integration of diagnosis and treatment",
]
tb = add_tb(slide4, Emu(700000), Emu(1700000), Emu(10500000), Emu(1100000))
tf = tb.text_frame
for j, point in enumerate(thesis_points):
    if j == 0:
        r = tf.paragraphs[0].add_run()
    else:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
    set_run(r, f"\u2022  {point}", 12, BODY)

# Section 2: INVESTOR FIT - teal header bar
add_rect(slide4, Emu(700000), Emu(3000000), Emu(10500000), Emu(400000), TEAL)
tb = add_tb(slide4, Emu(800000), Emu(3010000), Emu(10300000), Emu(380000))
r = tb.text_frame.paragraphs[0].add_run()
set_run(r, "INVESTOR FIT", 16, WHITE, True)

investors = [
    "Beauty/wellness funds — True Beauty Ventures, Imaginary Ventures, Silas Capital",
    "Health equity investors — Backstage Capital, Harlem Capital, Precursor Ventures",
    "Accelerators — Sephora Accelerate, YC, NAACP x L\u2019Or\u00e9al",
    "Domain angels — dermatologists, Toni Ko, skin founders",
]
tb = add_tb(slide4, Emu(700000), Emu(3500000), Emu(10500000), Emu(1000000))
tf = tb.text_frame
for j, point in enumerate(investors):
    if j == 0:
        r = tf.paragraphs[0].add_run()
    else:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
    set_run(r, f"\u2022  {point}", 12, BODY)

# Section 3: PRE-PITCH VALIDATION - amber warning bar style
add_rect(slide4, Emu(700000), Emu(4700000), Emu(10500000), Emu(1800000), RGBColor(0xFE, 0xF3, 0xC7))
tb = add_tb(slide4, Emu(800000), Emu(4750000), Emu(10300000), Emu(350000))
r = tb.text_frame.paragraphs[0].add_run()
set_run(r, "PRE-PITCH VALIDATION", 16, AMBER, True)

tb = add_tb(slide4, Emu(800000), Emu(5150000), Emu(10300000), Emu(1200000))
tf = tb.text_frame
r = tf.paragraphs[0].add_run()
set_run(r, "Run 50\u2013100 telehealth consults before pitching.", 14, BODY, True)
add_para(tf, "Cost: ~$4.5K for 100 consults ($2K platform + $25/consult)", 12, DKGRAY, space_before=6)
add_para(tf, "Proves: condition breakdown, diagnosis-to-purchase conversion rate", 12, DKGRAY, space_before=4)
p = add_para(tf, "", 13, RUST, True, space_before=12)
r = p.add_run()
set_run(r, "\u201CWe ran 100 consults, 78 converted to pre-orders, here\u2019s the condition breakdown, and now we know exactly what to build.\u201D", 13, RUST, True)


# ================================================================
# RECOLOR SLIDE 6 (Flywheel) - brand palette
# ================================================================
print("Recoloring Slide 6 (Flywheel)...")
slide6 = prs.slides[5]

brand_node_colors = {
    'top': (CREAM, RUST),        # Telehealth Consult - warm cream + rust text
    'right': (NAVY, WHITE),       # Condition Data - navy + white text
    'bottom': (RGBColor(0xE8, 0xF5, 0xE9), GREEN),  # Matched Product - light green + green text
    'left': (BLUSH, RED),         # Reorder + Retention - blush + red text
}

for shape in slide6.shapes:
    name = shape.name
    if not hasattr(shape, 'text'):
        continue

    txt = shape.text.strip()

    # Recolor rounded rectangles
    if 'Rounded Rectangle' in name:
        try:
            cy = shape.top + shape.height // 2
            cx = shape.left + shape.width // 2
            if cy < 2500000:
                shape.fill.solid(); shape.fill.fore_color.rgb = brand_node_colors['top'][0]
                shape.line.color.rgb = RUST; shape.line.width = Pt(1.5)
            elif cx > 7000000:
                shape.fill.solid(); shape.fill.fore_color.rgb = brand_node_colors['right'][0]
                shape.line.fill.background()
            elif cy > 4500000:
                shape.fill.solid(); shape.fill.fore_color.rgb = brand_node_colors['bottom'][0]
                shape.line.color.rgb = GREEN; shape.line.width = Pt(1.5)
            else:
                shape.fill.solid(); shape.fill.fore_color.rgb = brand_node_colors['left'][0]
                shape.line.color.rgb = RED; shape.line.width = Pt(1.5)
        except:
            pass

    # Recolor text in nodes
    if 'TELEHEALTH' in txt or 'CONSULT' in txt:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RUST
    elif 'CONDITION' in txt or 'DATA' in txt:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE
    elif 'MATCHED' in txt or 'PRODUCT' in txt:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = GREEN
    elif 'REORDER' in txt or 'RETENTION' in txt:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RED

    # Fix detail text colors
    if 'Provider confirms' in txt or '$5 margin' in txt:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = BODY
    elif 'Seb derm vs' in txt:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RGBColor(0xCA, 0xD5, 0xFC)  # Light on navy
    elif '$28 Steady' in txt or '75%' in txt:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = BODY
    elif 'Chronic condition' in txt:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = BODY

    # "THE LOOP" center
    if 'THE' in txt and 'LOOP' in txt:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE

    # Bottom insight text
    if 'Each turn' in txt or 'compounds' in txt:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.font.bold:
                    r.font.color.rgb = NAVY
                else:
                    r.font.color.rgb = BODY

# Add background tint
slide6.background.fill.solid()
slide6.background.fill.fore_color.rgb = RGBColor(0xFD, 0xFA, 0xF5)  # Warm off-white


# ================================================================
# RECOLOR SLIDE 7 (Curology) - brand cream/navy
# ================================================================
print("Recoloring Slide 7 (Curology)...")
slide7 = prs.slides[6]

for shape in slide7.shapes:
    if not hasattr(shape, 'fill'):
        continue
    name = shape.name

    # Curology column background - keep light gray
    if 'Rounded Rectangle 2' in name:
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            shape.line.color.rgb = GRAY
            shape.line.width = Pt(0.5)
        except:
            pass

    # Kyn Skyn column background - change from lavender to brand cream
    if 'Rounded Rectangle 5' in name:
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = CREAM
            shape.line.color.rgb = RUST
            shape.line.width = Pt(1.5)
        except:
            pass

# Fix Kyn Skyn header color
for shape in slide7.shapes:
    if hasattr(shape, 'text') and shape.text.strip() == 'KYN SKYN':
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RED

# Fix "Why Curology can't..." to use amber warning style
for shape in slide7.shapes:
    if hasattr(shape, 'text') and 'Why Curology' in shape.text:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.font.bold:
                    r.font.color.rgb = AMBER
                else:
                    r.font.color.rgb = BODY


# ================================================================
# FIX ALL FOOTNOTES to 8pt minimum
# ================================================================
print("Fixing footnotes to 8pt minimum...")
footnote_count = 0
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is not None and run.font.size < Pt(8):
                    run.font.size = Pt(8)
                    footnote_count += 1
print(f"  Fixed {footnote_count} runs below 8pt")


# ================================================================
# SAVE
# ================================================================
prs.save('C:/GitHub/Kyn-skyn/output/Kyn_skyn_final.pptx')
print("\nDesign rebuild complete.")
