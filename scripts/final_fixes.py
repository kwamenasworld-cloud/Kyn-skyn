"""Final fixes for all critical and high-priority issues."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

prs = Presentation('C:/GitHub/Kyn-skyn/output/Kyn_skyn_final.pptx')
ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'

NAVY = RGBColor(0x1E, 0x27, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
CREAM = RGBColor(0xFB, 0xF0, 0xED)
WARM_BROWN = RGBColor(0x8B, 0x6F, 0x5E)

# ================================================================
# FIX 1: SLIDE 13 - Garbled KPI strip text
# The "$559K+" overlaps "cumulative revenue" visually.
# Fix: ensure proper spacing and that text boxes don't overlap.
# ================================================================
print("FIX 1: Slide 13 KPI strip...")
slide13 = prs.slides[12]
for shape in slide13.shapes:
    if hasattr(shape, 'text'):
        # Fix the $559K+ text box - make sure it's not too tall
        if shape.text.strip() == '$559K+':
            # Reduce height and ensure it stays in its lane
            shape.height = Emu(250000)  # ~0.27 inches
            print(f"  Resized $559K+ box height to {shape.height}")

        # Fix "cumulative revenue" text - ensure it's positioned below $559K+
        if 'cumulative revenue' in shape.text:
            shape.top = Emu(5650000)  # Push down slightly
            shape.height = Emu(400000)
            # Also make text smaller to fit
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size > Pt(9):
                        run.font.size = Pt(9)
            print(f"  Repositioned 'cumulative revenue' box")

        # Fix the "2,000+" text if present
        if shape.text.strip() == '2,000+':
            shape.height = Emu(250000)

        if 'customers with' in shape.text:
            shape.top = Emu(5650000)
            shape.height = Emu(400000)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size > Pt(9):
                        run.font.size = Pt(9)

        if shape.text.strip() == 'CM+':
            shape.height = Emu(250000)

        if 'contribution' in shape.text and 'margin' in shape.text:
            shape.top = Emu(5650000)
            shape.height = Emu(400000)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size > Pt(9):
                        run.font.size = Pt(9)

print("  Done")

# ================================================================
# FIX 2: SLIDE 4 - Reframe to investor-facing
# Change meta-commentary headers to value proposition headers
# ================================================================
print("FIX 2: Slide 4 - Reframe for investors...")
slide4 = prs.slides[3]
for shape in slide4.shapes:
    if not hasattr(shape, 'text'):
        continue

    # Title
    if shape.text.strip() == 'Investability | Path from thesis to proof':
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = "Why this is investable | "
        r1.font.size = Pt(20)
        r1.font.bold = True
        r1.font.color.rgb = NAVY
        r2 = p.add_run()
        r2.text = "Vertical integration with a clear validation path"
        r2.font.size = Pt(20)
        r2.font.bold = True
        r2.font.color.rgb = DARK_TEXT
        print("  Reframed title")

    # "WHAT WORKS FOR VCs" -> "THE THESIS"
    if shape.text.strip() == 'WHAT WORKS FOR VCs':
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = "THE THESIS"
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)
        print("  Reframed 'What Works' -> 'The Thesis'")

    # Update the bullets underneath
    if 'Telehealth + product loop is a genuine moat' in shape.text:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "Telehealth + product loop creates a data moat"
        r.font.size = Pt(12)
        r.font.color.rgb = DARK_TEXT

        for text, bold in [
            ("Diagnosis feeds product, product data feeds roadmap", False),
            ("Clean unit economics: $14 contribution margin per unit", False),
            ("$250K ask on $1.5M cap — 18 months to breakeven", False),
            ("Not just skincare, not just telehealth — vertical integration", False),
        ]:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(4)
            r2 = p2.add_run()
            r2.text = text
            r2.font.size = Pt(12)
            r2.font.color.rgb = DARK_TEXT
            r2.font.bold = bold
        print("  Updated thesis bullets")

    # "TARGET INVESTOR PROFILE" -> "INVESTOR FIT"
    if shape.text.strip() == 'TARGET INVESTOR PROFILE':
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = "INVESTOR FIT"
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = NAVY
        print("  Reframed 'Target Investor Profile' -> 'Investor Fit'")

    # Update investor fit bullets
    if 'Beauty/wellness funds' in shape.text:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "Beauty/wellness (True Beauty Ventures, Imaginary)"
        r.font.size = Pt(12)
        r.font.color.rgb = DARK_TEXT

        for text in [
            "Health equity (Backstage, Harlem Capital, Precursor)",
            "Accelerators (Sephora Accelerate, YC)",
            "Domain angels (dermatologists, skin founders)",
        ]:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(4)
            r2 = p2.add_run()
            r2.text = text
            r2.font.size = Pt(12)
            r2.font.color.rgb = DARK_TEXT
        print("  Updated investor fit bullets")

    # "THE VALIDATION EXPERIMENT..." -> "PRE-PITCH VALIDATION"
    if 'VALIDATION EXPERIMENT' in shape.text:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = "PRE-PITCH VALIDATION"
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = NAVY
        print("  Reframed validation header")

    # Update validation content
    if 'Run 50' in shape.text and 'consults' in shape.text:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "Run 50\u2013100 telehealth consults before pitching."
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = DARK_TEXT

        for text in [
            "Cost: ~$4.5K for 100 consults ($2K platform + $25/consult)",
            "Proves: condition breakdown, diagnosis-to-purchase conversion rate",
        ]:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(6)
            r2 = p2.add_run()
            r2.text = text
            r2.font.size = Pt(12)
            r2.font.color.rgb = DARK_TEXT

        # Add the quote as final paragraph
        p3 = tf.add_paragraph()
        p3.space_before = Pt(12)
        r3 = p3.add_run()
        r3.text = ("\u201CWe ran 100 consults, 78 converted to pre-orders, here\u2019s the condition "
                   "breakdown, and now we know exactly what to build.\u201D")
        r3.font.size = Pt(13)
        r3.font.bold = True
        r3.font.color.rgb = WARM_BROWN
        print("  Updated validation content")

print("  Done")

# ================================================================
# FIX 3: SLIDE 19 - Closing slide contrast
# Make tagline bigger, ensure it's truly white, add a dark overlay
# ================================================================
print("FIX 3: Slide 19 - Closing contrast...")
slide19 = prs.slides[18]

# Find the tagline shape and force white with XML manipulation
for t_elem in slide19._element.iter(f'{{{ns_a}}}t'):
    if t_elem.text and 'diagnosis-to-treatment' in t_elem.text:
        # Get the parent <a:r> element
        r_elem = t_elem.getparent()
        rPr = r_elem.find(f'{{{ns_a}}}rPr')
        if rPr is None:
            rPr = etree.SubElement(r_elem, f'{{{ns_a}}}rPr')
            r_elem.insert(0, rPr)

        # Clear all existing fills
        for fill in rPr.findall(f'{{{ns_a}}}solidFill'):
            rPr.remove(fill)

        # Set white fill
        fill = etree.SubElement(rPr, f'{{{ns_a}}}solidFill')
        clr = etree.SubElement(fill, f'{{{ns_a}}}srgbClr')
        clr.set('val', 'FFFFFF')

        # Set size and bold
        rPr.set('sz', '2800')  # 28pt
        rPr.set('b', '1')

        # Also remove any schemeClr that might override
        for sc in rPr.findall(f'.//{{{ns_a}}}schemeClr'):
            sc.getparent().remove(sc)

        print("  Set tagline to white 28pt bold (deep XML fix)")

# Also find the background shape and try to darken it
# The slide likely has a background fill - let's check
bg = slide19.background
if bg.fill.type is not None:
    print(f"  Background fill type: {bg.fill.type}")
else:
    print("  No direct background fill found, checking shapes...")

# Add a semi-transparent dark overlay behind the tagline
# Actually, let's just add a dark rectangle behind the tagline text
tagline_shape = None
for shape in slide19.shapes:
    if hasattr(shape, 'text') and 'diagnosis-to-treatment' in shape.text:
        tagline_shape = shape
        break

if tagline_shape:
    # Add a dark navy rectangle behind the tagline
    overlay = slide19.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        tagline_shape.left - Emu(50000),
        tagline_shape.top - Emu(30000),
        tagline_shape.width + Emu(100000),
        tagline_shape.height + Emu(60000)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = NAVY
    overlay.line.fill.background()

    # Move overlay behind the text by reordering XML
    spTree = slide19.shapes._spTree
    overlay_el = overlay._element
    tagline_el = tagline_shape._element
    # Move overlay before tagline in the tree
    spTree.remove(overlay_el)
    idx = list(spTree).index(tagline_el)
    spTree.insert(idx, overlay_el)
    print("  Added dark navy rectangle behind tagline for contrast")

print("  Done")

# ================================================================
# FIX 4: SLIDE 6 - Make flywheel arrows bolder
# ================================================================
print("FIX 4: Slide 6 - Bolder flywheel arrows...")
slide6 = prs.slides[5]
cxns = slide6._element.findall(f'.//{{{ns_p}}}cxnSp')
for cxn in cxns:
    ln = cxn.find(f'.//{{{ns_a}}}ln')
    if ln is not None:
        # Increase line width from 38100 (3pt) to 76200 (6pt)
        ln.set('w', '76200')
        # Make the arrowhead larger
        tail = ln.find(f'{{{ns_a}}}tailEnd')
        if tail is not None:
            tail.set('w', 'lg')
            tail.set('len', 'lg')
        print(f"  Thickened arrow to 6pt")
print("  Done")

# ================================================================
# FIX 5: SLIDE 27 - Add Kyn Skyn position marker
# ================================================================
print("FIX 5: Slide 27 - Add Kyn Skyn position...")
slide27 = prs.slides[26]

# Add a text callout showing Kyn Skyn's position
# Position it at the right side of the slide as an overlay
callout = slide27.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Emu(9500000), Emu(1200000),  # Right side
    Emu(2400000), Emu(2000000)
)
callout.fill.solid()
callout.fill.fore_color.rgb = NAVY
callout.line.fill.background()

# Add Kyn Skyn position text
tb = slide27.shapes.add_textbox(
    Emu(9600000), Emu(1300000),
    Emu(2200000), Emu(1800000)
)
tf = tb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "KYN SKYN"
r.font.size = Pt(16)
r.font.bold = True
r.font.color.rgb = WHITE

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(8)
r2 = p2.add_run()
r2.text = "Sits at the intersection of all three"
r2.font.size = Pt(11)
r2.font.color.rgb = WHITE

p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(12)
r3 = p3.add_run()
r3.text = "\u2713 Clinical science\n\u2713 Condition-specific\n\u2713 Culturally resonant"
r3.font.size = Pt(11)
r3.font.color.rgb = RGBColor(0xE8, 0xBB, 0xB3)  # Coral accent

print("  Added Kyn Skyn position callout")
print("  Done")

# ================================================================
# SAVE
# ================================================================
prs.save('C:/GitHub/Kyn-skyn/output/Kyn_skyn_final.pptx')
print("\nAll fixes saved.")
