import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

prs = Presentation('C:/GitHub/Kyn-skyn/output/Kyn_skyn_updated.pptx')
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

print(f'Slide dimensions: {SLIDE_W/914400:.2f}" x {SLIDE_H/914400:.2f}"')
print(f'Total slides: {len(prs.slides)}')

for i, slide in enumerate(prs.slides):
    issues = []
    for shape in slide.shapes:
        left, top = shape.left, shape.top
        w, h = shape.width, shape.height

        # Off-slide elements
        if left + w > SLIDE_W + 50000:
            issues.append(f'OFF-SLIDE RIGHT: {shape.name} extends to {(left+w)/914400:.2f}" (slide width {SLIDE_W/914400:.2f}")')
        if top + h > SLIDE_H + 50000:
            issues.append(f'OFF-SLIDE BOTTOM: {shape.name} extends to {(top+h)/914400:.2f}" (slide height {SLIDE_H/914400:.2f}")')
        if left < -10000:
            issues.append(f'OFF-SLIDE LEFT: {shape.name} at x={left/914400:.2f}"')
        if top < -10000:
            issues.append(f'OFF-SLIDE TOP: {shape.name} at y={top/914400:.2f}"')

        # Zero-size elements with text
        if (w == 0 or h == 0) and hasattr(shape, 'text') and shape.text.strip():
            issues.append(f"ZERO-SIZE WITH TEXT: {shape.name} ({w}x{h}) text='{shape.text[:30]}'")

        # Very small elements that might be artifacts
        if w < 5000 and h < 5000 and w > 0 and h > 0:
            txt = shape.text[:20] if hasattr(shape, 'text') and shape.text else ''
            issues.append(f"TINY ELEMENT: {shape.name} ({w/914400:.3f}\" x {h/914400:.3f}\") text='{txt}'")

        # Elements at exactly (0,0) - potential misplaced elements
        if left == 0 and top == 0:
            txt = shape.text[:30] if hasattr(shape, 'text') and shape.text else ''
            issues.append(f'AT ORIGIN (0,0): {shape.name} ({w/914400:.2f}" x {h/914400:.2f}") text="{txt}"')

    # Check for duplicate shape names
    names = [s.name for s in slide.shapes]
    seen = set()
    for name in names:
        if name in seen:
            issues.append(f"DUPLICATE SHAPE NAME: '{name}'")
        seen.add(name)

    if issues:
        print(f'\nSLIDE {i+1}:')
        for issue in issues:
            print(f'  {issue}')

print('\n=== CHECK COMPLETE ===')
