"""Fix all text overflow and arrow visibility issues."""
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from lxml import etree

prs = Presentation('C:/GitHub/Kyn-skyn/output/Kyn_skyn_final.pptx')
ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
SW = prs.slide_width
SH = prs.slide_height

# ================================================================
# FIX 1: SLIDE 6 - Rebuild flywheel arrows as proper diagonal connectors
# ================================================================
print("FIX 1: Slide 6 flywheel arrows...")
slide6 = prs.slides[5]

# Remove old broken connectors
cxns = slide6._element.findall(f'.//{{{ns_p}}}cxnSp')
for cxn in cxns:
    cxn.getparent().remove(cxn)
print(f"  Removed {len(cxns)} broken connectors")

# Node box positions (from inspection):
# TOP (Consult): Rounded Rect at ~(4890000, 1200000) size 2400000x1200000
# RIGHT (Data): Rounded Rect at ~(7096000, 3200000) size 2400000x1200000
# BOTTOM (Product): Rounded Rect at ~(4890000, 4800000) size 2400000x1200000
# LEFT (Reorder): Rounded Rect at ~(2697000, 3200000) size 2400000x1200000

# Get actual positions from the rounded rectangles
rects = {}
for shape in slide6.shapes:
    if 'Rounded Rectangle' in shape.name:
        cy = shape.top + shape.height // 2
        cx = shape.left + shape.width // 2
        # Classify by position
        if cy < 2500000:
            rects['top'] = shape
        elif cx > 7000000:
            rects['right'] = shape
        elif cy > 4500000:
            rects['bottom'] = shape
        else:
            rects['left'] = shape

for name, s in rects.items():
    print(f"  {name}: L={s.left} T={s.top} W={s.width} H={s.height}")

def add_connector(slide, x1, y1, x2, y2, color_hex='1E2761', width=114300):
    """Add a properly-sized connector arrow."""
    spTree = slide.shapes._spTree

    # Ensure non-zero dimensions
    left = min(x1, x2)
    top = min(y1, y2)
    cx = max(abs(x2 - x1), 100)  # Minimum 100 EMU
    cy = max(abs(y2 - y1), 100)

    flipH = '1' if x2 < x1 else '0'
    flipV = '1' if y2 < y1 else '0'

    xml = f'''<p:cxnSp xmlns:a="{ns_a}" xmlns:p="{ns_p}" xmlns:r="{ns_r}">
      <p:nvCxnSpPr>
        <p:cNvPr id="0" name="Arrow"/>
        <p:cNvCxnSpPr/>
        <p:nvPr/>
      </p:nvCxnSpPr>
      <p:spPr>
        <a:xfrm flipH="{flipH}" flipV="{flipV}">
          <a:off x="{left}" y="{top}"/>
          <a:ext cx="{cx}" cy="{cy}"/>
        </a:xfrm>
        <a:prstGeom prst="straightConnector1">
          <a:avLst/>
        </a:prstGeom>
        <a:ln w="{width}" cap="rnd">
          <a:solidFill>
            <a:srgbClr val="{color_hex}"/>
          </a:solidFill>
          <a:tailEnd type="triangle" w="lg" len="lg"/>
        </a:ln>
      </p:spPr>
    </p:cxnSp>'''
    spTree.append(etree.fromstring(xml))

# Arrow 1: TOP → RIGHT (bottom-right of top box to top-left of right box)
t = rects['top']
r = rects['right']
add_connector(slide6,
    t.left + t.width,          t.top + t.height,       # bottom-right of top
    r.left,                     r.top)                   # top-left of right
print("  Added arrow: TOP → RIGHT")

# Arrow 2: RIGHT → BOTTOM (bottom-left of right box to top-right of bottom box)
b = rects['bottom']
add_connector(slide6,
    r.left,                     r.top + r.height,       # bottom-left of right
    b.left + b.width,          b.top)                   # top-right of bottom
print("  Added arrow: RIGHT → BOTTOM")

# Arrow 3: BOTTOM → LEFT (bottom-left of bottom box to bottom-right of left box)
l = rects['left']
add_connector(slide6,
    b.left,                     b.top,                   # top-left of bottom
    l.left + l.width,          l.top + l.height)        # bottom-right of left
print("  Added arrow: BOTTOM → LEFT")

# Arrow 4: LEFT → TOP (top-right of left box to bottom-left of top box)
add_connector(slide6,
    l.left + l.width,          l.top,                   # top-right of left
    t.left,                     t.top + t.height)       # bottom-left of top
print("  Added arrow: LEFT → TOP")


# ================================================================
# FIX 2: SLIDE 2 - Widen/fix the right treatment box text
# The text boxes extend slightly past their container
# ================================================================
print("\nFIX 2: Slide 2 treatment text overflow...")
slide2 = prs.slides[1]

# Find the right container boundary
right_container = None
for shape in slide2.shapes:
    if shape.name == 'Rectangle: Rounded Corners 33':
        right_container = shape
        break

if right_container:
    rc_right = right_container.left + right_container.width
    print(f"  Right container edge: {rc_right}")

    # Fix text boxes that extend past the container
    for shape in slide2.shapes:
        if hasattr(shape, 'text') and shape.text:
            shape_right = shape.left + shape.width
            # If text box extends past container, shrink it
            if shape.left > right_container.left and shape_right > rc_right:
                new_width = rc_right - shape.left - Emu(50000)
                if new_width > Emu(100000):  # Don't make boxes too small
                    shape.width = new_width
                    print(f"  Shrunk {shape.name}: \"{shape.text[:40]}\"")
                else:
                    print(f"  Skipped {shape.name} (would be too narrow)")

# Also fix the "Consumers looking for solutions" text that extends past slide
for shape in slide2.shapes:
    if hasattr(shape, 'text') and 'Consumers are looking' in shape.text:
        shape_right = shape.left + shape.width
        if shape_right > SW:
            shape.width = SW - shape.left - Emu(100000)
            print(f"  Shrunk consumers text to fit slide")


# ================================================================
# FIX 3: SLIDE 5 - Make cycle text boxes smaller font to prevent overflow
# ================================================================
print("\nFIX 3: Slide 5 cycle text overflow...")
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if hasattr(shape, 'text'):
        t = shape.text.strip()
        # These are the cycle step descriptions
        if t.startswith('Consumer applies') or t.startswith('Malassezia fungus') or \
           t.startswith('Fungal overgrowth') or t.startswith('Consumer buys'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size > Pt(10):
                        run.font.size = Pt(10)
                        print(f"  Reduced font: \"{t[:40]}\"")
                    elif not run.font.size:
                        run.font.size = Pt(10)
                        print(f"  Set font 10pt: \"{t[:40]}\"")


# ================================================================
# FIX 4: SLIDE 13 - Fix footnote collision with R&D line
# ================================================================
print("\nFIX 4: Slide 13 footnote collision...")
slide13 = prs.slides[12]
for shape in slide13.shapes:
    if hasattr(shape, 'text') and 'Marketing remains' in shape.text:
        # Move the marketing note text further up
        shape.top = Emu(5000000)
        shape.height = Emu(280000)
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(7)
        print(f"  Moved marketing note up, reduced to 7pt")


# ================================================================
# FIX 5: SLIDE 15 - Ensure Phase 3 bullets don't clip
# Reduce font size in Phase 3 bullet text
# ================================================================
print("\nFIX 5: Slide 15 GTM phase text...")
slide15 = prs.slides[14]
# Find shapes related to Phase 3 content
for shape in slide15.shapes:
    if hasattr(shape, 'text') and 'Prepare for launch' in shape.text:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size and run.font.size > Pt(9):
                    run.font.size = Pt(9)
        print(f"  Reduced Phase 3 text size")


# ================================================================
# SAVE
# ================================================================
prs.save('C:/GitHub/Kyn-skyn/output/Kyn_skyn_final.pptx')
print("\nAll overflow fixes saved.")
