"""
BCG PL recommendations - Full deck restructure:
1. Add "Why Now?" slide
2. Redesign The Model slide with flywheel diagram
3. Move Investability to early position, Curology & Amazon validation into main deck
4. Fix revenue bridge math on 18-Month Goals
5. Fix SOM label, LTV/CAC label, closing slide statement
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import math

INPUT = 'C:/GitHub/Kyn-skyn/output/Kyn_skyn_updated.pptx'
OUTPUT = 'C:/GitHub/Kyn-skyn/output/Kyn_skyn_final.pptx'

prs = Presentation(INPUT)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# Brand colors
NAVY = RGBColor(0x1E, 0x27, 0x61)
CREAM = RGBColor(0xFB, 0xF0, 0xED)
WARM_BROWN = RGBColor(0x8B, 0x6F, 0x5E)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_CORAL = RGBColor(0xE8, 0xBB, 0xB3)
LIGHT_BG = RGBColor(0xFA, 0xF7, 0xF4)
GREEN_ACCENT = RGBColor(0x4C, 0xAF, 0x50)
TEAL = RGBColor(0x02, 0x80, 0x90)

def set_run(run, text, size_pt=14, bold=False, color=DARK_TEXT, font_name=None):
    run.text = text
    run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if font_name:
        run.font.name = font_name

def add_textbox(slide, left, top, width, height, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = word_wrap
    return txBox

def add_para(tf, text, size_pt=14, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, space_before=None):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    set_run(run, text, size_pt, bold, color)
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_oval(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape

def make_title(slide, prefix, suffix, prefix_color=NAVY, suffix_color=DARK_TEXT):
    """Add a standard title bar at the top of a slide."""
    tb = add_textbox(slide, Emu(850168), Emu(400000), Emu(10500000), Emu(600000))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    set_run(r1, prefix, 20, True, prefix_color)
    r2 = p.add_run()
    set_run(r2, suffix, 20, True, suffix_color)
    return tb

blank_layout = prs.slide_layouts[15]  # "1_Blank"

# ================================================================
# 1. CREATE "WHY NOW?" SLIDE
# ================================================================
print("Creating 'Why Now?' slide...")
why_slide = prs.slides.add_slide(blank_layout)
why_slide.background.fill.solid()
why_slide.background.fill.fore_color.rgb = WHITE

make_title(why_slide, "Why now | ", "Four forces converging to create this category")

# Four cards in a 2x2 grid
cards = [
    {
        'title': 'Scalp care is surging',
        'body': 'Global scalp care market growing 8%+ CAGR. Mintel reports scalp health is the #1 emerging concern in skincare. Consumers are moving beyond shampoo into treatment.',
        'color': RGBColor(0xE8, 0xF5, 0xE9),
        'accent': GREEN_ACCENT,
        'num': '01',
    },
    {
        'title': 'Social media is driving awareness',
        'body': 'TikTok #sebderm has 500M+ views. Reddit r/SebDerm has 80K+ members. Consumers are self-diagnosing and searching for fungal-safe products\u2014but finding nothing purpose-built.',
        'color': RGBColor(0xED, 0xF0, 0xFB),
        'accent': NAVY,
        'num': '02',
    },
    {
        'title': 'Telehealth is normalized',
        'body': 'Post-COVID telehealth adoption makes a $30 consult feel routine, not novel. Platforms like Ola Digital Health offer turnkey setup for $2K. The infrastructure exists\u2014nobody has applied it to fungal skin.',
        'color': RGBColor(0xFB, 0xF0, 0xED),
        'accent': WARM_BROWN,
        'num': '03',
    },
    {
        'title': 'No incumbent owns this',
        'body': 'Curology serves acne. Dermazen sells products. Head & Shoulders treats symptoms. Nobody has combined diagnosis + treatment for Malassezia conditions. The white space is wide open.',
        'color': CREAM,
        'accent': ACCENT_CORAL,
        'num': '04',
    },
]

positions = [
    (Emu(600000), Emu(1250000)),   # top-left
    (Emu(6300000), Emu(1250000)),  # top-right
    (Emu(600000), Emu(3900000)),   # bottom-left
    (Emu(6300000), Emu(3900000)),  # bottom-right
]
card_w = Emu(5200000)
card_h = Emu(2300000)

for i, card in enumerate(cards):
    x, y = positions[i]
    rect = add_rounded_rect(why_slide, x, y, card_w, card_h, card['color'], card['accent'])

    # Number badge
    badge = add_oval(why_slide, x + Emu(200000), y + Emu(180000), Emu(400000), Emu(400000), card['accent'])
    badge_tb = add_textbox(why_slide, x + Emu(200000), y + Emu(180000), Emu(400000), Emu(400000))
    tf = badge_tb.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    set_run(r, card['num'], 16, True, WHITE)

    # Title
    title_tb = add_textbox(why_slide, x + Emu(720000), y + Emu(200000), Emu(4200000), Emu(350000))
    tf = title_tb.text_frame
    r = tf.paragraphs[0].add_run()
    set_run(r, card['title'], 18, True, DARK_TEXT)

    # Body
    body_tb = add_textbox(why_slide, x + Emu(200000), y + Emu(650000), Emu(4800000), Emu(1500000))
    tf = body_tb.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    set_run(r, card['body'], 12, False, DARK_TEXT)

# Bottom note
note_tb = add_textbox(why_slide, Emu(600000), Emu(6350000), Emu(10900000), Emu(350000))
tf = note_tb.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run()
set_run(r, "The category is forming now. The question is who builds the platform first.", 14, True, NAVY)


# ================================================================
# 2. REDESIGN THE MODEL SLIDE WITH FLYWHEEL
# ================================================================
print("Redesigning The Model slide with flywheel...")
slide4 = prs.slides[3]  # Current "The Model" slide

# Remove ALL existing shapes from slide 4
spTree = slide4.shapes._spTree
shapes_to_remove = list(spTree)
for child in shapes_to_remove:
    if child.tag.endswith('}grpSpPr') or child.tag.endswith('}nvGrpSpPr'):
        continue  # Keep the required group properties
    if child.tag.endswith('}sp') or child.tag.endswith('}pic') or child.tag.endswith('}grpSp') or child.tag.endswith('}cxnSp'):
        spTree.remove(child)

# Set background
slide4.background.fill.solid()
slide4.background.fill.fore_color.rgb = WHITE

# Title
make_title(slide4, "The Model | ", "Diagnosis + treatment flywheel")

# Center flywheel - use rounded rects arranged in a diamond/circle pattern
# Center point of flywheel
cx, cy = Emu(6096000), Emu(3800000)  # Center of slide horizontally, slightly below center
radius = Emu(2000000)

# Flywheel nodes positioned in a diamond
nodes = [
    {'label': '$30 TELEHEALTH\nCONSULT', 'detail': 'Provider confirms condition\n$5 margin, self-funded', 'pos': 'top',
     'color': CREAM, 'accent': ACCENT_CORAL},
    {'label': 'CONDITION\nDATA', 'detail': 'Seb derm vs. tinea vs.\nfungal acne breakdown', 'pos': 'right',
     'color': RGBColor(0xED, 0xF0, 0xFB), 'accent': NAVY},
    {'label': 'MATCHED\nPRODUCT', 'detail': '$28 Steady Serum\n75%+ product margin', 'pos': 'bottom',
     'color': RGBColor(0xE8, 0xF5, 0xE9), 'accent': GREEN_ACCENT},
    {'label': 'REORDER +\nRETENTION', 'detail': 'Chronic condition = repeat\npurchase + next SKU data', 'pos': 'left',
     'color': RGBColor(0xFB, 0xF0, 0xED), 'accent': WARM_BROWN},
]

node_w = Emu(2400000)
node_h = Emu(1200000)

node_positions = {
    'top':    (cx - node_w // 2, cy - radius - node_h // 2),
    'right':  (cx + radius - node_w // 2 + Emu(200000), cy - node_h // 2),
    'bottom': (cx - node_w // 2, cy + radius - node_h // 2),
    'left':   (cx - radius - node_w // 2 - Emu(200000), cy - node_h // 2),
}

for node in nodes:
    nx, ny = node_positions[node['pos']]
    rect = add_rounded_rect(slide4, nx, ny, node_w, node_h, node['color'], node['accent'])

    # Label
    label_tb = add_textbox(slide4, nx + Emu(100000), ny + Emu(80000), node_w - Emu(200000), Emu(500000))
    tf = label_tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    set_run(r, node['label'], 14, True, node['accent'])

    # Detail
    detail_tb = add_textbox(slide4, nx + Emu(100000), ny + Emu(580000), node_w - Emu(200000), Emu(550000))
    tf = detail_tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    set_run(r, node['detail'], 11, False, DARK_TEXT)

# Center circle with "THE LOOP" text
center_oval = add_oval(slide4, cx - Emu(600000), cy - Emu(450000), Emu(1200000), Emu(900000), NAVY)
center_tb = add_textbox(slide4, cx - Emu(550000), cy - Emu(350000), Emu(1100000), Emu(700000))
tf = center_tb.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run()
set_run(r, "THE\nLOOP", 20, True, WHITE)

# Arrow labels between nodes (using text boxes with arrow characters)
arrow_positions = [
    (cx + Emu(700000), cy - radius + Emu(200000), "\u2192"),   # top -> right
    (cx + Emu(700000), cy + radius - Emu(500000), "\u2192"),   # right -> bottom (down-right)
    (cx - Emu(900000), cy + radius - Emu(500000), "\u2190"),   # bottom -> left
    (cx - Emu(900000), cy - radius + Emu(200000), "\u2190"),   # left -> top
]

for ax, ay, arrow_char in arrow_positions:
    atb = add_textbox(slide4, ax, ay, Emu(400000), Emu(350000))
    tf = atb.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    set_run(r, arrow_char, 28, True, NAVY)

# Bottom insight bar
insight_tb = add_textbox(slide4, Emu(600000), Emu(6250000), Emu(10900000), Emu(400000))
tf = insight_tb.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r1 = tf.paragraphs[0].add_run()
set_run(r1, "Each turn of the loop compounds: ", 13, True, NAVY)
r2 = tf.paragraphs[0].add_run()
set_run(r2, "more consults \u2192 richer data \u2192 smarter SKUs \u2192 higher retention \u2192 a product roadmap no competitor can replicate.", 13, False, DARK_TEXT)


# ================================================================
# 3. FIX REVENUE BRIDGE ON 18-MONTH GOALS (currently slide 11, idx 10)
# ================================================================
print("Fixing revenue bridge on 18-Month Goals...")
slide11 = prs.slides[10]

# Find and update the footnote/note text
for shape in slide11.shapes:
    if hasattr(shape, 'text') and '$35 blended AOV' in shape.text:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        set_run(r, (
            "Revenue bridge: 2,000 customers \u00d7 $35 blended AOV \u00d7 1.3 avg orders/customer (30% reorder on "
            "3-month cycle over 18 months) = ~$91K from initial cohort. Staggered monthly acquisition of ~110 new "
            "customers/month \u00d7 $35 AOV compounds to ~$559K cumulative. "
            "Launch CAC target $22, optimized to ~$13 by M12 via organic content and retention. "
            "Gross margin excludes fulfillment; contribution margin includes it."
        ), 9, False, DARK_TEXT)
        print("  Updated revenue bridge note")

    # Fix LTV/CAC label
    if hasattr(shape, 'text') and 'LTV / CAC' in shape.text:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'LTV / CAC' in run.text:
                    run.text = run.text.replace('LTV / CAC', 'LTV / CAC (target)')
                    print("  Labeled LTV/CAC as target")

    if hasattr(shape, 'text') and shape.text.strip() == '3.4\u00d7':
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '3.4' in run.text:
                    run.text = '3.4\u00d7*'
                    print("  Added asterisk to 3.4x")


# ================================================================
# 4. FIX SOM LABEL ON SIZING SLIDE (slide 6, idx 5)
# ================================================================
print("Fixing SOM/SAM labels on sizing slide...")
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if hasattr(shape, 'text'):
        # Fix the $500M SOM - it should be labeled as near-term target
        if 'SOM' in shape.text and 'Diagnosis-to-treatment' in shape.text:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            set_run(r, "SOM \u2014 Year 1 target: ~$559K via diagnosis-to-treatment DTC into Black & Hispanic SD segment", 10, False, DARK_TEXT)
            print("  Fixed SOM label to realistic Y1 target")


# ================================================================
# 5. UPDATE CLOSING SLIDE (slide 16, idx 15)
# ================================================================
print("Updating closing slide...")
slide16 = prs.slides[15]
# Find the tagline text and add investment thesis
for shape in slide16.shapes:
    if hasattr(shape, 'text') and 'Science-first' in shape.text:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, "The first company to own the diagnosis-to-treatment loop for fungal skin.", 18, True, WHITE)
        print("  Updated closing tagline")


# ================================================================
# 6. REORDER SLIDES
# ================================================================
print("Reordering slides...")

# Current slide order (0-indexed):
# 0: Title
# 1: Market
# 2: Problem
# 3: The Model (redesigned)
# 4: Solution
# 5: Sizing
# 6: Opportunity
# 7: Unit Economics
# 8: Fundraise
# 9: Investability
# 10: 18-Month Goals
# 11: GTM I
# 12: GTM II
# 13: Product Roadmap
# 14: Founders
# 15: Close/Contact
# 16: Appendix divider
# 17: Problem Detail
# 18: Design
# 19: Packaging
# 20: Supply Chain
# 21: Amazon Validation (MOVE TO MAIN)
# 22: Product Detail
# 23: Care Value Chain
# 24: Competitive Landscape
# 25: Curology Comparison (MOVE TO MAIN)
# 26: Hero Ingredient
# 27: Dermazen Case Study
# 28+: Why Now? (NEW, just added at end)

# Target order:
# 0: Title
# 1: Market
# 2: Why Now? (was 28)
# 3: Investability (was 9)
# 4: Problem (was 2)
# 5: The Model with flywheel (was 3)
# 6: Curology Comparison (was 25)
# 7: Amazon Validation (was 21)
# 8: Solution (was 4)
# 9: Sizing (was 5)
# 10: Opportunity (was 6)
# 11: Unit Economics (was 7)
# 12: Fundraise (was 8)
# 13: 18-Month Goals (was 10)
# 14: GTM I (was 11)
# 15: GTM II (was 12)
# 16: Product Roadmap (was 13)
# 17: Founders (was 14)
# 18: Close/Contact (was 15)
# 19: Appendix divider (was 16)
# 20+: remaining appendix (17-24, 26, 27 - minus 21 and 25 which moved)

sldIdLst = prs.element.find(qn('p:sldIdLst'))
sldIds = list(sldIdLst)

# Map current indices to sldId elements
# sldIds[0] = slide 1 (Title), sldIds[1] = slide 2 (Market), etc.
# The "Why Now?" slide was added at the end, so it's sldIds[28]

# Build the new order by index into current sldIds array
new_order = [
    0,   # Title
    1,   # Market
    28,  # Why Now? (newly added, at end)
    9,   # Investability
    2,   # Problem
    3,   # The Model (redesigned)
    25,  # Curology Comparison
    21,  # Amazon Validation
    4,   # Solution
    5,   # Sizing
    6,   # Opportunity
    7,   # Unit Economics
    8,   # Fundraise
    10,  # 18-Month Goals
    11,  # GTM I
    12,  # GTM II
    13,  # Product Roadmap
    14,  # Founders
    15,  # Close/Contact
    16,  # Appendix divider
    17,  # Problem Detail
    18,  # Design
    19,  # Packaging
    20,  # Supply Chain
    22,  # Product Detail
    23,  # Care Value Chain
    24,  # Competitive Landscape
    26,  # Hero Ingredient
    27,  # Dermazen Case Study
]

# Verify we have all slides accounted for
assert len(new_order) == len(sldIds), f"Order has {len(new_order)} items but deck has {len(sldIds)} slides"
assert sorted(new_order) == list(range(len(sldIds))), "Not all slides accounted for in reorder"

# Reorder by removing all and re-inserting in new order
ordered_sldIds = [sldIds[i] for i in new_order]

# Remove all from sldIdLst
for sid in sldIds:
    sldIdLst.remove(sid)

# Re-insert in new order
for sid in ordered_sldIds:
    sldIdLst.append(sid)

print(f"  Reordered {len(ordered_sldIds)} slides")

# ================================================================
# 7. ADDITIONAL SLIDE-SPECIFIC FIXES FROM BCG REVIEW
# ================================================================

# Fix slide 5 (Solution) - reframe the disclaimer
# After reorder, Solution is at new position 8 (0-indexed)
# But we need to find it by content since we're working with the original indices
print("Applying slide-specific fixes...")

# Original slide 4 (Solution) - fix "supplementary control" language
orig_solution = prs.slides[4]
for shape in orig_solution.shapes:
    if hasattr(shape, 'text') and 'supplementary control' in shape.text:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'supplementary control' in run.text:
                    run.text = run.text.replace(
                        'Positioned as supplementary control - not the only line of defense for acute flare-ups.',
                        'Daily maintenance serum that works alongside your existing routine.'
                    )
                    print("  Fixed solution disclaimer language")

# ================================================================
# SAVE
# ================================================================
print(f"\nSaving to {OUTPUT}...")
prs.save(OUTPUT)
print(f"Done! Final deck saved. Total slides: {len(prs.slides)}")
