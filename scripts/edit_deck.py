"""
Comprehensive deck edit script:
1. Update existing slides with "diagnosis + treatment" framing
2. Add new slides: The Model, Curology Differentiation, VC Investability
3. Fix slide 16 artifact (off-slide rectangle)
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy
from lxml import etree
import os

INPUT = 'C:/Users/admin/OneDrive/Documents/Kyn_skyn.pptx'
OUTPUT = 'C:/GitHub/Kyn-skyn/output/Kyn_skyn_updated.pptx'

os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)

prs = Presentation(INPUT)
SLIDE_W = prs.slide_width   # 12192000 EMU
SLIDE_H = prs.slide_height  # 6858000 EMU

# Color palette from brand
NAVY = RGBColor(0x1E, 0x27, 0x61)
CREAM = RGBColor(0xFB, 0xF0, 0xED)
WARM_BROWN = RGBColor(0x8B, 0x6F, 0x5E)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_CORAL = RGBColor(0xE8, 0xBB, 0xB3)
LIGHT_BG = RGBColor(0xFA, 0xF7, 0xF4)

def set_text_props(run, size_pt=None, bold=None, color=None, font_name=None):
    if size_pt:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if font_name:
        run.font.name = font_name

def clear_and_set(text_frame, text, size_pt=14, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_text_props(run, size_pt, bold, color)
    return p

def add_paragraph(text_frame, text, size_pt=14, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, space_before=None):
    p = text_frame.add_paragraph()
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    set_text_props(run, size_pt, bold, color)
    return p

def add_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

# ============================================================
# 1. UPDATE SLIDE 1 - Title slide
# ============================================================
print("Editing Slide 1 - Title...")
slide1 = prs.slides[0]
for shape in slide1.shapes:
    # Update subtitle
    if shape.name == 'TextBox 2':
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "Diagnosis-driven skincare for fungal skin conditions"
        set_text_props(run, 18, False, DARK_TEXT)
    # Update thesis - Text 3 inside Group 35
    if shape.name == 'Group 35':
        for child in shape.shapes:
            if child.name == 'Text 3':
                tf = child.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = ("Kyn Skyn is building the first diagnosis-to-treatment platform "
                           "for Malassezia-driven skin conditions. A $30 telehealth consultation "
                           "confirms which fungal condition the consumer has (seb derm, tinea versicolor, "
                           "or fungal acne). Then we match them with fungal-safe products formulated "
                           "specifically for what the provider just confirmed.")
                set_text_props(run, 12, False, DARK_TEXT)
            if child.name == 'Text 4':
                tf = child.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = ("The diagnosis makes the treatment sale stronger. The treatment data makes "
                           "the next product smarter. Two things, one customer \u2014 that\u2019s the loop.")
                set_text_props(run, 12, False, DARK_TEXT)

# ============================================================
# 2. UPDATE SLIDE 4 - Our Solution (reframe to include diagnosis)
# ============================================================
print("Editing Slide 4 - Solution...")
slide4 = prs.slides[3]
for shape in slide4.shapes:
    if shape.name == 'Text 3':
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Our solution | "
        set_text_props(run, 20, True, None)
        run2 = p.add_run()
        run2.text = "Diagnosis confirms the condition, Steady Serum treats it \u2014 fungal-safe by design"
        set_text_props(run2, 20, True, None)

# ============================================================
# 3. UPDATE SLIDE 7 - Unit Economics (add telehealth line)
# ============================================================
print("Editing Slide 7 - Unit Economics...")
slide7 = prs.slides[6]
for shape in slide7.shapes:
    if shape.name == 'Title 2':
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Unit Economics | "
        set_text_props(run, 20, True, None)
        run2 = p.add_run()
        run2.text = "Product margin funds growth; telehealth consult is self-funded at $5 margin per visit"
        set_text_props(run2, 20, True, None)
    if shape.name == 'Text 58':
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = ("Note: Product unit economics at $28 price. Telehealth consult: $30 price, ~$25 cost "
                   "(Ola Digital Health platform), ~$5 margin per consult. Consult is self-funded and "
                   "serves as the acquisition tool. 75%+ product margin is the primary revenue driver.")
        set_text_props(run, 9, False, DARK_TEXT)

# ============================================================
# 4. FIX SLIDE 16 ARTIFACT (remove off-slide rectangle)
# ============================================================
print("Fixing Slide 16 artifact...")
slide16 = prs.slides[15]
shapes_to_remove = []
for shape in slide16.shapes:
    if shape.name == 'Rounded Rectangle 132':
        # This shape is at x=-6096, extends off-slide
        shapes_to_remove.append(shape)

for shape in shapes_to_remove:
    sp = shape._element
    sp.getparent().remove(sp)
    print(f"  Removed: {shape.name}")

# ============================================================
# 5. ADD NEW SLIDE: "The Model" (after slide 3, before solution)
# ============================================================
print("Adding new slide: The Model...")

# We'll duplicate the blank layout and build the slide
blank_layout = prs.slide_layouts[15]  # "1_Blank"

# Insert after slide 3 (index 3) - we need to manipulate XML for ordering
# For now, add at end and reorder later
model_slide = prs.slides.add_slide(blank_layout)

# Background
bg = model_slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = WHITE

# Title
title_box = add_textbox(model_slide, Emu(850168), Emu(400000), Emu(10500000), Emu(650000))
tf = title_box.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "The Model | "
set_text_props(run, 20, True, NAVY)
run2 = p.add_run()
run2.text = "Two things, one customer \u2014 diagnosis drives treatment, treatment data drives the roadmap"
set_text_props(run2, 20, True, DARK_TEXT)

# Left column - Diagnosis
diag_rect = add_rounded_rect(model_slide, Emu(600000), Emu(1400000), Emu(5200000), Emu(4600000), CREAM, ACCENT_CORAL)
diag_rect.text_frame.word_wrap = True

# Diagnosis header
diag_header = add_textbox(model_slide, Emu(800000), Emu(1550000), Emu(4800000), Emu(500000))
tf = diag_header.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "DIAGNOSIS"
set_text_props(run, 24, True, NAVY)

# Diagnosis content
diag_content = add_textbox(model_slide, Emu(800000), Emu(2100000), Emu(4800000), Emu(3600000))
tf = diag_content.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
run = p.add_run()
run.text = "$30 telehealth consultation"
set_text_props(run, 16, True, DARK_TEXT)

add_paragraph(tf, "Confirms which Malassezia condition: seb derm, tinea versicolor, or fungal acne", 13, False, DARK_TEXT, space_before=8)
add_paragraph(tf, "", 8)
add_paragraph(tf, "Self-funded model", 14, True, DARK_TEXT, space_before=12)
add_paragraph(tf, "$30 consult price", 13, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "~$25 cost (Ola Digital Health)", 13, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "~$5 margin per consult", 13, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "", 8)
add_paragraph(tf, "Platform: $2K setup + $25/consult", 13, False, WARM_BROWN, space_before=8)

# Right column - Treatment
treat_rect = add_rounded_rect(model_slide, Emu(6300000), Emu(1400000), Emu(5200000), Emu(4600000), RGBColor(0xED, 0xF0, 0xFB), NAVY)

# Treatment header
treat_header = add_textbox(model_slide, Emu(6500000), Emu(1550000), Emu(4800000), Emu(500000))
tf = treat_header.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "TREATMENT"
set_text_props(run, 24, True, NAVY)

# Treatment content
treat_content = add_textbox(model_slide, Emu(6500000), Emu(2100000), Emu(4800000), Emu(3600000))
tf = treat_content.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
run = p.add_run()
run.text = "Fungal-safe products matched to diagnosis"
set_text_props(run, 16, True, DARK_TEXT)

add_paragraph(tf, "Starting with Steady Serum ($28, 75%+ margin), expanding to 7-SKU platform", 13, False, DARK_TEXT, space_before=8)
add_paragraph(tf, "", 8)
add_paragraph(tf, "Product economics", 14, True, DARK_TEXT, space_before=12)
add_paragraph(tf, "$28 unit price", 13, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "$5.50 COGS", 13, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "~$14 contribution margin", 13, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "", 8)
add_paragraph(tf, "75%+ product margin is the revenue engine", 13, False, WARM_BROWN, space_before=8)

# Bottom - The Loop
loop_box = add_textbox(model_slide, Emu(600000), Emu(6150000), Emu(10900000), Emu(500000))
tf = loop_box.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "The loop: "
set_text_props(run, 13, True, NAVY)
run2 = p.add_run()
run2.text = ("Diagnosis makes the treatment sale stronger. Treatment data makes the next product smarter. "
            "Condition breakdown from consults drives the product roadmap.")
set_text_props(run2, 13, False, DARK_TEXT)


# ============================================================
# 6. ADD NEW SLIDE: Curology Differentiation
# ============================================================
print("Adding new slide: Curology Differentiation...")

curology_slide = prs.slides.add_slide(blank_layout)
bg = curology_slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = WHITE

# Title
title_box = add_textbox(curology_slide, Emu(850168), Emu(400000), Emu(10500000), Emu(650000))
tf = title_box.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Curology comparison | "
set_text_props(run, 20, True, NAVY)
run2 = p.add_run()
run2.text = "Curology validated telehealth-to-skincare. We apply the model to fungal conditions they can\u2019t serve"
set_text_props(run2, 20, True, DARK_TEXT)

# Curology column
cur_rect = add_rounded_rect(curology_slide, Emu(600000), Emu(1350000), Emu(5200000), Emu(3400000), RGBColor(0xF5, 0xF5, 0xF5), RGBColor(0xCC, 0xCC, 0xCC))

cur_header = add_textbox(curology_slide, Emu(800000), Emu(1450000), Emu(4800000), Emu(400000))
tf = cur_header.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "CUROLOGY"
set_text_props(run, 20, True, RGBColor(0x66, 0x66, 0x66))

cur_content = add_textbox(curology_slide, Emu(800000), Emu(1900000), Emu(4800000), Emu(2700000))
tf = cur_content.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
run = p.add_run()
run.text = "Custom prescription compounding"
set_text_props(run, 14, True, DARK_TEXT)

add_paragraph(tf, "Every customer gets a unique formula mixed at their pharmacy", 12, False, DARK_TEXT, space_before=6)
add_paragraph(tf, "", 6)
add_paragraph(tf, "Serves acne (saturated market)", 12, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "Requires pharmacy infrastructure", 12, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "$20\u201340/month subscription model", 12, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "Higher per-unit costs (custom formulation)", 12, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "$100M+ raised, millions of customers", 12, False, DARK_TEXT, space_before=4)

# Kyn Skyn column
kyn_rect = add_rounded_rect(curology_slide, Emu(6300000), Emu(1350000), Emu(5200000), Emu(3400000), CREAM, NAVY)

kyn_header = add_textbox(curology_slide, Emu(6500000), Emu(1450000), Emu(4800000), Emu(400000))
tf = kyn_header.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "KYN SKYN"
set_text_props(run, 20, True, NAVY)

kyn_content = add_textbox(curology_slide, Emu(6500000), Emu(1900000), Emu(4800000), Emu(2700000))
tf = kyn_content.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
run = p.add_run()
run.text = "Standard diagnosis \u2192 standard SKU"
set_text_props(run, 14, True, DARK_TEXT)

add_paragraph(tf, "Everyone with seb derm gets the same proven Steady Serum", 12, False, DARK_TEXT, space_before=6)
add_paragraph(tf, "", 6)
add_paragraph(tf, "Serves fungal conditions (underserved)", 12, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "No pharmacy needed (OTC/cosmetic)", 12, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "$28 one-time + reorder model", 12, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "75%+ margin (same bottle for everyone)", 12, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "$2K platform + $25/consult (lean)", 12, False, DARK_TEXT, space_before=4)

# Why this matters box
why_box = add_textbox(curology_slide, Emu(600000), Emu(5000000), Emu(10900000), Emu(700000))
tf = why_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Why Curology can\u2019t easily add Malassezia: "
set_text_props(run, 13, True, NAVY)
run2 = p.add_run()
run2.text = ("Curology is a prescription compounding business. Fungal conditions are treated with "
            "OTC cosmetics and OTC antifungals, not custom compounds. They\u2019d have to build a "
            "separate non-prescription product line \u2014 a different business entirely.")
set_text_props(run2, 13, False, DARK_TEXT)

# Pitch line
pitch_box = add_textbox(curology_slide, Emu(600000), Emu(5900000), Emu(10900000), Emu(600000))
tf = pitch_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = ("\u201CCurology proved telehealth-to-skincare works. They did it for acne with custom prescriptions. "
           "We\u2019re doing it for fungal conditions with targeted OTC products. "
           "Their model requires a pharmacy. Ours requires a $2K platform and a $28 bottle.\u201D")
set_text_props(run, 13, True, WARM_BROWN)


# ============================================================
# 7. ADD NEW SLIDE: Path to VC Investable
# ============================================================
print("Adding new slide: VC Investability...")

vc_slide = prs.slides.add_slide(blank_layout)
bg = vc_slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = WHITE

# Title
title_box = add_textbox(vc_slide, Emu(850168), Emu(400000), Emu(10500000), Emu(650000))
tf = title_box.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Investability | "
set_text_props(run, 20, True, NAVY)
run2 = p.add_run()
run2.text = "What makes this venture-backable and the validation that turns thesis into proof"
set_text_props(run2, 20, True, DARK_TEXT)

# Left column - What works
works_rect = add_rounded_rect(vc_slide, Emu(600000), Emu(1350000), Emu(5200000), Emu(2200000), RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0x4C, 0xAF, 0x50))

works_header = add_textbox(vc_slide, Emu(800000), Emu(1450000), Emu(4800000), Emu(350000))
tf = works_header.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "WHAT WORKS FOR VCs"
set_text_props(run, 16, True, RGBColor(0x2E, 0x7D, 0x32))

works_content = add_textbox(vc_slide, Emu(800000), Emu(1850000), Emu(4800000), Emu(1600000))
tf = works_content.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Telehealth + product loop is a genuine moat"
set_text_props(run, 12, False, DARK_TEXT)
add_paragraph(tf, "Diagnosis feeds product, product data feeds roadmap", 12, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "Unit economics are clean ($14 contribution margin)", 12, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "$250K ask on $1.5M cap signals discipline", 12, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "Not just a skincare brand, not just telehealth", 12, False, DARK_TEXT, space_before=4)

# Right column - Target investors
inv_rect = add_rounded_rect(vc_slide, Emu(6300000), Emu(1350000), Emu(5200000), Emu(2200000), RGBColor(0xED, 0xF0, 0xFB), NAVY)

inv_header = add_textbox(vc_slide, Emu(6500000), Emu(1450000), Emu(4800000), Emu(350000))
tf = inv_header.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "TARGET INVESTOR PROFILE"
set_text_props(run, 16, True, NAVY)

inv_content = add_textbox(vc_slide, Emu(6500000), Emu(1850000), Emu(4800000), Emu(1600000))
tf = inv_content.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Beauty/wellness funds (True Beauty Ventures, Imaginary)"
set_text_props(run, 12, False, DARK_TEXT)
add_paragraph(tf, "Health equity investors (Backstage, Harlem Capital, Precursor)", 12, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "Accelerators (NAACP x L\u2019Or\u00e9al, Sephora Accelerate, YC)", 12, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "Domain angels (Toni Ko, dermatologists, skin founders)", 12, False, DARK_TEXT, space_before=4)

# Bottom - Validation experiment
val_rect = add_rounded_rect(vc_slide, Emu(600000), Emu(3850000), Emu(10900000), Emu(2700000), CREAM, ACCENT_CORAL)

val_header = add_textbox(vc_slide, Emu(800000), Emu(3950000), Emu(10500000), Emu(350000))
tf = val_header.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "THE VALIDATION EXPERIMENT THAT MAKES THIS CLEARLY INVESTABLE"
set_text_props(run, 16, True, NAVY)

val_content = add_textbox(vc_slide, Emu(800000), Emu(4400000), Emu(10500000), Emu(2000000))
tf = val_content.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
run = p.add_run()
run.text = "Run 50\u2013100 telehealth consults before pitching."
set_text_props(run, 14, True, DARK_TEXT)

add_paragraph(tf, "Cost: ~$2K platform setup + $25/consult = ~$4.5K for 100 consults", 12, False, DARK_TEXT, space_before=6)
add_paragraph(tf, "Show: condition breakdown (% seb derm vs. tinea vs. fungal acne)", 12, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "Show: conversion rate from diagnosis to waitlist/pre-order", 12, False, DARK_TEXT, space_before=4)
add_paragraph(tf, "", 8)

p2 = add_paragraph(tf, "", 14, True, WARM_BROWN, space_before=8)
run = p2.add_run()
run.text = ("\u201CWe ran 100 consults, 78 converted to pre-orders, here\u2019s the condition breakdown, "
           "and now we know exactly what to build.\u201D \u2014 That pitch is worth 10x more than a polished deck.")
set_text_props(run, 13, True, WARM_BROWN)


# ============================================================
# 8. REORDER SLIDES - move new slides to correct positions
# ============================================================
print("Reordering slides...")

# Current order: 25 original slides + 3 new ones at end (indices 25, 26, 27)
# New slides:
#   - "The Model" (index 25) -> should go after slide 3 (The Problem), before slide 4 (Our Solution)
#   - "Curology Differentiation" (index 26) -> should go after competitive landscape in appendix (slide 23)
#   - "VC Investability" (index 27) -> should go after fundraise slide (slide 8)

# We need to manipulate the XML sldIdLst
# Access the presentation XML
from pptx.oxml.ns import qn

pres_xml = prs.element
sldIdLst = pres_xml.find(qn('p:sldIdLst'))
sldIds = list(sldIdLst)

# Get the 3 new slide entries (last 3)
model_sldId = sldIds[-3]    # The Model
curology_sldId = sldIds[-2]  # Curology
vc_sldId = sldIds[-1]       # VC Investability

# Remove them from current positions
sldIdLst.remove(model_sldId)
sldIdLst.remove(curology_sldId)
sldIdLst.remove(vc_sldId)

# Re-read the list after removal
sldIds = list(sldIdLst)
# Now we have 25 entries (original slides)

# Insert "The Model" after slide 3 (index 2, 0-based) -> position 3
sldIdLst.insert(3, model_sldId)

# Re-read after insertion
sldIds = list(sldIdLst)
# Now 26 entries. Fundraise was slide 8 (index 7), now index 8 due to insertion

# Insert "VC Investability" after fundraise slide (original slide 8, now at index 8)
sldIdLst.insert(9, vc_sldId)

# Re-read after insertion
sldIds = list(sldIdLst)
# Now 27 entries. Competitive landscape was slide 23 (appendix). Original index 22.
# After two insertions before it: +2 = index 24
# Insert "Curology Differentiation" after it
sldIdLst.insert(25, curology_sldId)

print("Final slide order set.")

# ============================================================
# 9. SAVE
# ============================================================
print(f"Saving to {OUTPUT}...")
prs.save(OUTPUT)
print("Done! Deck saved successfully.")
print(f"Total slides: {len(prs.slides)}")
